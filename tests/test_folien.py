"""PPTX/PDF → PNG über LibreOffice und pdftoppm."""

import shutil

import pytest

from app import folien

hat_werkzeuge = pytest.mark.skipif(
    not folien.werkzeuge_vorhanden(),
    reason="LibreOffice/pdftoppm nicht installiert (läuft im Container)")


def test_werkzeuge_vorhanden_prueft_beide():
    erwartet = bool(shutil.which("soffice") or shutil.which("libreoffice")) \
        and bool(shutil.which("pdftoppm"))
    assert folien.werkzeuge_vorhanden() is erwartet


def test_fehlende_quelle_wirft(tmp_path):
    with pytest.raises(folien.FolienFehler, match="nicht gefunden"):
        folien.exportiere(tmp_path / "gibt-es-nicht.pptx", tmp_path)


def test_unbekanntes_format_wirft(tmp_path):
    quelle = tmp_path / "notiz.txt"
    quelle.write_text("x")
    with pytest.raises(folien.FolienFehler, match="Format"):
        folien.exportiere(quelle, tmp_path)


@hat_werkzeuge
def test_pptx_wird_zu_pngs(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    quelle = tmp_path / "deck.pptx"
    p = Presentation()
    for text in ("Erste Folie", "Zweite Folie"):
        folie = p.slides.add_slide(p.slide_layouts[5])
        folie.shapes.title.text = text
    p.save(quelle)

    ziel = tmp_path / "folien"
    bilder = folien.exportiere(quelle, ziel)
    assert [b.name for b in bilder] == ["folie-01.png", "folie-02.png"]
    assert all(b.stat().st_size > 0 for b in bilder)


@hat_werkzeuge
def test_zielordner_wird_vorher_geleert(tmp_path):
    from pptx import Presentation

    quelle = tmp_path / "deck.pptx"
    p = Presentation()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(quelle)

    ziel = tmp_path / "folien"
    ziel.mkdir()
    (ziel / "folie-99.png").write_bytes(b"alt")
    bilder = folien.exportiere(quelle, ziel)
    assert not (ziel / "folie-99.png").exists()
    assert len(bilder) == 1


def test_nummerierung_bleibt_sortierbar_ueber_99_folien(tmp_path):
    """pdftoppm padded je nach Gesamtseitenzahl (3-stellig ab 100 Seiten).

    Ohne Rücksicht darauf würde ein fixes zweistelliges Format
    "folie-100.png" alphabetisch vor "folie-11.png" einsortieren. Diese
    Normalisierung läuft ohne LibreOffice/pdftoppm — sie simuliert nur deren
    Ausgabe für ein 100-Folien-Deck.
    """
    ziel = tmp_path / "folien"
    ziel.mkdir()
    for n in range(1, 101):
        (ziel / f"folie-{n:03d}.png").write_bytes(b"x")

    folien._normalisiere_nummerierung(ziel)

    namen = [b.name for b in sorted(ziel.glob("folie-*.png"))]
    assert namen[0] == "folie-001.png"
    assert namen[9] == "folie-010.png"
    assert namen[10] == "folie-011.png"
    assert namen[-1] == "folie-100.png"
    assert namen == sorted(namen)


def test_fehlschlag_laesst_kein_teilergebnis_in_ziel_dir(tmp_path, monkeypatch):
    """Bricht pdftoppm mitten im Deck ab, darf ziel_dir nicht halbfertig sein.

    Simuliert eine kaputte Seite spät im Deck: ein paar PNGs entstehen,
    dann schlägt der Aufruf fehl. ziel_dir muss danach entweder gar nicht
    existieren oder exakt seinen alten Inhalt zeigen — nie einen Rest aus
    dem gescheiterten Versuch.
    """
    quelle = tmp_path / "deck.pptx"
    quelle.write_text("x")
    monkeypatch.setattr(folien, "werkzeuge_vorhanden", lambda: True)
    monkeypatch.setattr(folien, "_als_pdf", lambda q, arbeit: arbeit / "deck.pdf")

    def kaputt(pdf, ziel_bilder, dpi):
        (ziel_bilder / "folie-01.png").write_bytes(b"halb")
        raise folien.FolienFehler("pdftoppm fehlgeschlagen: kaputte Seite")

    monkeypatch.setattr(folien, "_als_pngs", kaputt)

    # Fall 1: ziel_dir existiert noch nicht.
    ziel = tmp_path / "folien"
    with pytest.raises(folien.FolienFehler):
        folien.exportiere(quelle, ziel)
    assert not ziel.exists()

    # Fall 2: ziel_dir hat bereits ein gültiges früheres Ergebnis.
    ziel.mkdir()
    (ziel / "folie-01.png").write_bytes(b"alt-aber-gut")
    with pytest.raises(folien.FolienFehler):
        folien.exportiere(quelle, ziel)
    assert (ziel / "folie-01.png").read_bytes() == b"alt-aber-gut"
    assert [p.name for p in ziel.glob("*")] == ["folie-01.png"]
